import os
import sys
import re
from pptx import Presentation

# Ajout du chemin pour trouver le dossier shared
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from shared.database import SessionLocal
from shared.models import Gamme, Product

def extract_sections(text):
    """Extrait intelligemment les parties du texte via Regex."""
    sections = {
        "indications": "Non spécifié",
        "composition": "Non spécifié",
        "conseils": "Non spécifié"
    }

    # On utilise re.IGNORECASE pour être plus souple
    # On cherche le contenu entre les mots-clés
    try:
        # Cherche entre INDICATIONS et (COMPOSITION ou CONSEILS ou Fin)
        ind_match = re.search(r"INDICATIONS(.*?)(COMPOSITION|CONSEILS|$)", text, re.S | re.I)
        if ind_match:
            sections["indications"] = ind_match.group(1).strip()
    except: pass

    try:
        # Cherche entre COMPOSITION et (CONSEILS ou Fin)
        comp_match = re.search(r"COMPOSITION(.*?)(CONSEILS|$)", text, re.S | re.I)
        if comp_match:
            sections["composition"] = comp_match.group(1).strip()
    except: pass

    try:
        # Cherche après CONSEILS jusqu'à la fin
        cons_match = re.search(r"CONSEILS(.*?)$", text, re.S | re.I)
        if cons_match:
            sections["conseils"] = cons_match.group(1).strip()
    except: pass

    return sections

def import_from_pptx(pptx_folder):
    db = SessionLocal()
    
    if not os.path.exists(pptx_folder):
        print(f"Erreur : Le dossier {pptx_folder} n'existe pas.")
        return

    # On récupère les gammes pour le matching
    gammes = db.query(Gamme).all()
    gamme_map = {g.name.lower(): g.id for g in gammes}

    count_added = 0
    count_skipped = 0

    for filename in os.listdir(pptx_folder):
        if filename.endswith(".pptx"):
            file_path = os.path.join(pptx_folder, filename)
            clean_filename = filename.lower().replace("gamme", "").replace(".pptx", "").strip()
            
            # 1. Matching de la gamme
            found_gamme_id = None
            for g_name, g_id in gamme_map.items():
                if g_name in clean_filename or clean_filename in g_name:
                    found_gamme_id = g_id
                    break
            
            if not found_gamme_id:
                found_gamme_id = gamme_map.get("vital") # Fallback
                if not found_gamme_id: continue

            print(f"📖 Analyse intelligente de {filename}...")
            
            try:
                prs = Presentation(file_path)
                # On utilise enumerate pour sauter la 1ère slide sans faire de [1:] qui cause l'erreur rId
                for i, slide in enumerate(prs.slides):
                    if i == 0: 
                        continue # Saute la 1ère slide (Titre de la Gamme)

                    # Tri des shapes par position (pour avoir le titre en premier)
                    shapes = sorted([s for s in slide.shapes if hasattr(s, "text") and s.text.strip()], 
                                   key=lambda s: (s.top, s.left))
                    
                    if not shapes: 
                        continue
                    
                    # On cherche le premier nom de produit valide (pas une entête de section)
                    product_name = ""
                    for s in shapes:
                        txt = s.text.strip().split('\n')[0]
                        # On ignore les mots clés de section et les textes trop courts
                        if len(txt) > 3 and not any(kw in txt.upper() for kw in ["INDICATION", "COMPOSITION", "CONSEILS"]):
                            product_name = txt
                            break
                    
                    if not product_name: 
                        continue

                    # Vérifier les doublons
                    exists = db.query(Product).filter(Product.name == product_name).first()
                    if exists:
                        count_skipped += 1
                        continue
                    
                    # Fusionner tout le texte de la slide pour le Regex
                    full_text = "\n".join([s.text.strip() for s in shapes])
                    
                    # Extraction Regex
                    sections = extract_sections(full_text)
                    
                    # Construction de la nouvelle description structurée
                    final_description = f"--- INDICATIONS ---\n{sections['indications']}\n"
                    final_description += f"\n--- COMPOSITION ---\n{sections['composition']}\n"
                    final_description += f"\n--- CONSEILS D'UTILISATION ---\n{sections['conseils']}"

                    new_product = Product(
                        name=product_name,
                        gamme_id=found_gamme_id,
                        description=final_description
                    )
                    db.add(new_product)
                    count_added += 1
                    
            except Exception as e:
                print(f"❌ Erreur sur {filename}: {e}")

    db.commit()
    db.close()
    print(f"\n✨ Import intelligent terminé !")
    print(f"Produits ajoutés : {count_added}")
    print(f"Produits ignorés : {count_skipped}")

if __name__ == "__main__":
    dossier_ppts = "dso1/Data/PowerPoints"
    import_from_pptx(dossier_ppts)
